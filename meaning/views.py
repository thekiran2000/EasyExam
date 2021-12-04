from django.shortcuts import render

from math import ceil
import json
from PyDictionary import PyDictionary
from pptx import Presentation
import googletrans
import re 
import PyPDF2
from googletrans import Translator

import docx2txt

import pdfplumber

import gtts

from playsound import playsound

import os

from .models import Contact






# Create your views here.
from django.http import HttpResponse


str="""
"""




def index(request):
  
    return render(request, 'meaning/index.html')

def about(request):
    return render(request, 'meaning/about.html')

def contact(request):
   
    if request.method=="POST":
        name = request.POST.get('name', '')
        email = request.POST.get('email', '')
        phone = request.POST.get('phone', '')
        desc = request.POST.get('desc', '')
        contact = Contact(name=name, email=email, phone=phone, desc=desc)
        contact.save()
    
    
    return render(request, 'meaning/contact.html')

def notes(request):
    
    return render(request, 'meaning/notes.html')

def pro(request):
    if request.method=="POST":
        my_text=request.POST.get('pron','')
        my_code=request.POST.get('key','')
        print(my_text)
        print(my_code)
        # my_text=request.POST.get('','')
        tts = gtts.gTTS(my_text, lang=my_code)
        tts.save("hello.mp3")
        playsound("hello.mp3")
        os.remove("hello.mp3")


    return render(request, 'meaning/index.html')

def transfield(request):
    if request.method=="POST":
        my_text=request.POST.get('text','')
        lang=request.POST.get('language','')
        final={}
        translator = Translator(service_urls=['translate.googleapis.com'])
        result1=translator.translate(my_text,dest=lang)
        final[lang]=result1.text
        print(final[lang])

    return render(request,'meaning/index.html',{'final':final})

def result(request):
    if request.method=="POST":
        print(request)
        file=request.POST.get('doc','')
        file2=request.POST.get('ppt','')
        file3=request.POST.get('pdf','')
        file4=request.POST.get('hindi','')
        file5=request.POST.get('Marathi','')
        file6=request.POST.get('Comphindi','')
        file7=request.POST.get('CompMarathi','')
        number_set={'1','2','3','4','5','6','7','8','9','0'}

     
   
         
        
      
        if("docx" in file):

            my_text = docx2txt.process(f"D:/pythonreaddoc/{file}")
            res = re.sub(r'[^\w\s]', '', my_text)

            men_text=res.split()
        
            dictionary=PyDictionary()
            final={}
            for ele in men_text:
                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:
                    

                    final[ele]=dictionary.meaning(ele)
        
        elif("ppt" in file2):
            print("hiii")
            prs = Presentation(f"D:/pythonreaddoc/{file2}")

            string=""
            for slide in prs.slides:
                for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            string=string+shape.text
                            
                            
            res = re.sub(r'[^\w\s]', '', string)

            men_text=res.split()
        
            dictionary=PyDictionary()
            final={}
            for ele in men_text:

                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:

                        final[ele]=dictionary.meaning(ele)


        elif("pdf" in file3):
            string=""
            pdf = PyPDF2.PdfFileReader(open(f"D:/pythonreaddoc/{file3}", "rb"))

            for page in pdf.pages:
               
                string=string+page.extractText()
            


            res = re.sub(r'[^\w\s]', '', string)

            men_text=res.split()
        
            dictionary=PyDictionary()
            final={}
            for ele in men_text:

                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:

                        final[ele]=dictionary.meaning(ele)
        



        elif("docx" in file4):

            my_text = docx2txt.process(f"D:/pythonreaddoc/{file4}")
            res = re.sub(r'[^\w\s]', '', my_text)

            men_text=res.split()
        
            # dictionary=PyDictionary()
            
            translator = Translator(service_urls=['translate.googleapis.com'])

            final={}
            for ele in men_text :
                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:
                    result1=translator.translate(ele,dest='hi')
                    final[ele]=result1.text

        
        elif("ppt" in file4):
            print("hiii")
            prs = Presentation(f"D:/pythonreaddoc/{file4}")

            string=""
            for slide in prs.slides:
                for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            string=string+shape.text
                            
                            
            res = re.sub(r'[^\w\s]', '', string)

            men_text=res.split()
        
            translator = Translator()

            final={}
            for ele in men_text:
                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:
                    result1=translator.translate(ele,dest='hi')
                    final[ele]=result1.text


        elif("pdf" in file4):
            my_text=""
         
            final={}
            i=0
            with pdfplumber.open(f"D:/pythonreaddoc/{file4}") as pdf:
                for page in pdf.pages:
                    my_text=my_text+page.extract_text()
                    
                  
            


            res = re.sub(r'[^\w\s]', '', my_text)

            men_text=res.split()
        
            
            translator = Translator(service_urls=['translate.googleapis.com'])

            final={}
            for ele in men_text:
                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:
                    result1=translator.translate(ele,dest='hi')
                    final[ele]=result1.text


        elif("docx" in file5):

            my_text = docx2txt.process(f"D:/pythonreaddoc/{file5}")
            res = re.sub(r'[^\w\s]', '', my_text)

            men_text=res.split()
        
            
            translator = Translator(service_urls=['translate.googleapis.com'])

            final={}
            for ele in men_text:
                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:
                   result1=translator.translate(ele,dest='mr')
                   final[ele]=result1.text
        
        elif("ppt" in file5):
            print("hiii")
            prs = Presentation(f"D:/pythonreaddoc/{file5}")

            string=""
            for slide in prs.slides:
                for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            string=string+shape.text
                            
                            
            res = re.sub(r'[^\w\s]', '', string)

            men_text=res.split()
        
           
            translator = Translator(service_urls=['translate.googleapis.com'])

            final={}
            for ele in men_text:
                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:
                    result1=translator.translate(ele,dest='mr')
                    final[ele]=result1.text


        elif("pdf" in file5):
            
            my_text=""
         
            final={}
            i=0
            with pdfplumber.open(f"D:/pythonreaddoc/{file5}") as pdf:
                for page in pdf.pages:
                    my_text=my_text+page.extract_text()
                    
                  
            


            res = re.sub(r'[^\w\s]', '', my_text)

            men_text=res.split()
        
            
            translator = Translator(service_urls=['translate.googleapis.com'])

            final={}
            for ele in men_text:
                k=list(number_set.intersection(set(ele)))
                if(len(ele)>5) and len(k)==0:
                    result1=translator.translate(ele,dest='mr')
                    final[ele]=result1.text


        
        elif("docx" in file6):

            my_text = docx2txt.process(f"D:/pythonreaddoc/{file6}")
            # res = re.sub(r'[^\w\s]', '', my_text)

           
        
            # dictionary=PyDictionary()
            
            translator = Translator(service_urls=['translate.googleapis.com'])

            final={}
            
            result1=translator.translate(my_text,dest='hi')
            final["Translation"]=result1.text

        
        elif("ppt" in file6):
         
            prs = Presentation(f"D:/pythonreaddoc/{file6}")

            string=""
            final={}
            i=0
            translator = Translator(service_urls=['translate.googleapis.com'])
            for slide in prs.slides:
                for shape in slide.shapes:
                        i+=1
                        if hasattr(shape, "text"):
                            string=shape.text
                            
                            my_text=string
                            result=translator.translate(my_text, dest='mr')
                            final[i] =result.text


        elif("pdf" in file6):
            my_text=""
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            i=0
            with pdfplumber.open(f"D:/pythonreaddoc/{file6}") as pdf:
                for page in pdf.pages:
                    i+=1
                    my_text=page.extract_text()
                    result=translator.translate(my_text, dest='hi')
                    final[i] =result.text


        elif("docx" in file7):

            my_text = docx2txt.process(f"D:/pythonreaddoc/{file7}")
            
            # res = re.sub(r'[^\w\s]', '', my_text)

            # men_text=res.split()
        
            
            translator = Translator(service_urls=['translate.googleapis.com'])

            final={}
            
            result1=translator.translate(my_text,dest='mr')
            final["Translation"]=result1.text
        
        elif("ppt" in file7):
            print("hiii")
            prs = Presentation(f"D:/pythonreaddoc/{file7}")

            string=""
            final={}
            i=0
            translator = Translator(service_urls=['translate.googleapis.com'])
            for slide in prs.slides:
                for shape in slide.shapes:
                        i+=1
                        if hasattr(shape, "text"):
                            string=shape.text
                            
                            my_text=string
                            result=translator.translate(my_text, dest='mr')
                            final[i] =result.text
                            
            # res = re.sub(r'[^\w\s]', '', string)

            # men_text=res.split()
        
           
            

            
            
            # result1=translator.translate(string,dest='mr')
            # final["Translation"]=result1.text


        elif("pdf" in file7):
            my_text=""
            translator = Translator(service_urls=['translate.googleapis.com'])
            final={}
            i=0
            with pdfplumber.open(f"D:/pythonreaddoc/{file7}") as pdf:
                for page in pdf.pages:
                    i+=1
                    my_text=page.extract_text()
                    result=translator.translate(my_text, dest='mr')
                    final[i] =result.text
                    
    print(final)
  

    return render(request,'meaning/result.html',{'final':final})


